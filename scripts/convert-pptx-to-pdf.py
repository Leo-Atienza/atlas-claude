#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "comtypes>=1.4; sys_platform == 'win32'",
#     "python-pptx>=1.0",
#     "pymupdf>=1.24",
# ]
# ///
"""
convert-pptx-to-pdf.py — render a PowerPoint deck to PDF with full fidelity.

Unlike the wiki-extract-* scripts (which read text OUT of office files), this
DRIVES A REAL RENDERING ENGINE so the PDF looks exactly like the slides —
fonts, images, charts, layout, theme. python-pptx / PyMuPDF cannot render; they
are used here only for the --self-test fixture and page-count verification.

Engines (auto-detected, in preference order for --engine auto):
  1. Microsoft PowerPoint via COM    — perfect fidelity; honors "shrink text on
     overflow" autofit. Preferred when Office is installed.
  2. LibreOffice headless (soffice)  — fully headless (no PowerPoint session), but
     ignores autofit, so overflowing text boxes can CLIP on export. Free fallback.
     Install: winget install --id TheDocumentFoundation.LibreOffice

Usage:
  uv run convert-pptx-to-pdf.py <deck.pptx>                      # -> <deck.pdf> beside it
  uv run convert-pptx-to-pdf.py <deck.pptx> --output out.pdf     # explicit output (single input)
  uv run convert-pptx-to-pdf.py a.pptx b.pptx --outdir <dir>     # batch into a directory
  uv run convert-pptx-to-pdf.py <deck.pptx> --engine powerpoint  # force an engine
  uv run convert-pptx-to-pdf.py <deck.pptx> --json               # structured result
  uv run convert-pptx-to-pdf.py --self-test                      # synthesize+convert, exit 0/1

Flags:
  --output <file>   Output PDF path. Only valid with a single input.
  --outdir <dir>    Write PDFs into this directory (good for batch).
  --engine <name>   auto (default) | powerpoint | libreoffice
  --json            Emit {"results":[{input,output,engine,pages,bytes,status}]}.
  --self-test       Build a 2-slide deck, convert it, assert 2-page PDF, exit 0/1.
  -h, --help        Show this help.

Output (default): absolute path of each produced PDF, one per line (stdout).
Exit codes: 0 ok · 1 runtime error (incl. any failed conversion) · 2 usage error.

LIMITATIONS:
- Needs a rendering engine on the host (PowerPoint or LibreOffice). No engine -> exit 1.
- PowerPoint COM runs in the user's interactive session. If PowerPoint is already
  open with other decks, this only closes the deck it opened (never your work).
- First run resolves wheels via uv (~10-30s) and, for the COM path, generates the
  PowerPoint type-library wrapper once. Cached afterwards.
"""
import sys
import os
import glob
import json
import shutil
import tempfile
import subprocess
from pathlib import Path

# Force UTF-8 on stdout/stderr so non-ASCII paths/content don't crash on Windows
# where the default codec is cp1252. Python 3.7+; harmless on Linux/macOS.
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except (AttributeError, ValueError):
    pass

PP_SAVE_AS_PDF = 32          # PpSaveAsFileType.ppSaveAsPDF
PP_ALERTS_NONE = 1           # PpAlertLevel.ppAlertsNone
USAGE = (
    "usage: uv run convert-pptx-to-pdf.py <deck.pptx> [more.pptx ...] "
    "[--output <file> | --outdir <dir>] [--engine auto|powerpoint|libreoffice] "
    "[--json] | --self-test\n"
)
# Accepted PowerPoint input extensions — PowerPoint COM and LibreOffice both
# render all of these (.pptm/.ppsm are macro-enabled; .pps/.ppsx are shows;
# .pot*/.potx are templates).
PPT_EXTS = (".pptx", ".pptm", ".ppt", ".pps", ".ppsx", ".ppsm", ".potx", ".potm", ".pot")


# --------------------------------------------------------------------------- #
# Path handling — COM needs absolute Windows paths; args may arrive Git-Bash style
# --------------------------------------------------------------------------- #
def to_windows_abs(p):
    """Resolve to an absolute Windows path. Handles `~`, relative, and Git-Bash
    `/c/Users/...` drive paths (which native Windows Python doesn't understand)."""
    s = os.path.expanduser(str(p))
    # /c/Users/... or /c -> C:/Users/... (MSYS2/Git-Bash drive form)
    if len(s) >= 2 and s[0] == "/" and s[1].isalpha() and (len(s) == 2 or s[2] == "/"):
        s = s[1].upper() + ":" + s[2:]
    return os.path.abspath(s)


# --------------------------------------------------------------------------- #
# Engine detection
# --------------------------------------------------------------------------- #
def find_libreoffice():
    exe = shutil.which("soffice") or shutil.which("soffice.exe")
    if exe:
        return exe
    for p in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if os.path.isfile(p):
            return p
    return None


def find_powerpoint():
    for pat in (
        r"C:\Program Files\Microsoft Office\root\Office*\POWERPNT.EXE",
        r"C:\Program Files (x86)\Microsoft Office\root\Office*\POWERPNT.EXE",
        r"C:\Program Files\Microsoft Office\Office*\POWERPNT.EXE",
        r"C:\Program Files (x86)\Microsoft Office\Office*\POWERPNT.EXE",
    ):
        hits = glob.glob(pat)
        if hits:
            return hits[0]
    return None


def resolve_engine(choice):
    """Return (engine_name, engine_path). Raises RuntimeError if unavailable."""
    lo = find_libreoffice()
    pp = find_powerpoint()
    if choice == "libreoffice":
        if not lo:
            raise RuntimeError(
                "LibreOffice not found. Install it with:\n"
                "  winget install --id TheDocumentFoundation.LibreOffice"
            )
        return ("libreoffice", lo)
    if choice == "powerpoint":
        if not pp:
            raise RuntimeError("Microsoft PowerPoint (POWERPNT.EXE) not found on this machine.")
        return ("powerpoint", pp)
    if choice != "auto":
        raise RuntimeError(f"unknown engine: {choice!r} (use auto|powerpoint|libreoffice)")
    # auto: PowerPoint first when present — it honors "shrink text on overflow"
    # autofit, which LibreOffice ignores (LibreOffice can clip overflowing text
    # boxes on PDF export). Fall back to LibreOffice (fully headless) otherwise.
    if pp:
        return ("powerpoint", pp)
    if lo:
        return ("libreoffice", lo)
    raise RuntimeError(
        "No conversion engine found. Install one of:\n"
        "  Microsoft PowerPoint (Office), or\n"
        "  winget install --id TheDocumentFoundation.LibreOffice   (free, headless)."
    )


# --------------------------------------------------------------------------- #
# Converters
# --------------------------------------------------------------------------- #
def convert_powerpoint(in_abs, out_abs):
    """Drive Microsoft PowerPoint via COM. Closes only the deck it opens; quits
    the app only if no other presentations remain open."""
    import comtypes.client

    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    # NB: do NOT set Visible=False — PowerPoint (unlike Word/Excel) raises on it.
    try:
        powerpoint.DisplayAlerts = PP_ALERTS_NONE  # silence overwrite/repair dialogs
    except Exception:
        pass

    pres = None
    try:
        pres = powerpoint.Presentations.Open(in_abs, ReadOnly=True, Untitled=False, WithWindow=False)
        pres.SaveAs(out_abs, PP_SAVE_AS_PDF)
        pres.Close()
        pres = None
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            if powerpoint.Presentations.Count == 0:
                powerpoint.Quit()
        except Exception:
            pass
    return out_abs


def convert_libreoffice(soffice, in_abs, out_abs):
    """Convert via `soffice --headless --convert-to pdf`, in an isolated user
    profile so it works even when a LibreOffice GUI instance is already running."""
    out_dir = os.path.dirname(out_abs) or "."
    with tempfile.TemporaryDirectory(prefix="lo_profile_") as profile, \
         tempfile.TemporaryDirectory(prefix="lo_out_") as tmp_out:
        cmd = [
            soffice, "--headless", "--norestore", "--invisible", "--nologo",
            f"-env:UserInstallation={Path(profile).as_uri()}",
            "--convert-to", "pdf", "--outdir", tmp_out, in_abs,
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        produced = Path(tmp_out) / (Path(in_abs).stem + ".pdf")
        if proc.returncode != 0 or not produced.is_file():
            raise RuntimeError(
                f"LibreOffice conversion failed (rc={proc.returncode}). "
                f"stderr: {proc.stderr.strip() or '<empty>'}"
            )
        os.makedirs(out_dir, exist_ok=True)
        if os.path.exists(out_abs):
            os.remove(out_abs)
        shutil.move(str(produced), out_abs)
    return out_abs


def convert(engine_name, engine_path, in_abs, out_abs):
    if engine_name == "powerpoint":
        return convert_powerpoint(in_abs, out_abs)
    return convert_libreoffice(engine_path, in_abs, out_abs)


def pdf_page_count(pdf_path):
    """Best-effort page count via PyMuPDF; None if unavailable."""
    try:
        import fitz
        with fitz.open(pdf_path) as doc:
            return doc.page_count
    except Exception:
        return None


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #
def parse_args(argv):
    args = {"inputs": [], "output": None, "outdir": None, "engine": "auto",
            "json": False, "self_test": False, "help": False}
    i = 0
    while i < len(argv):
        a = argv[i]
        if a in ("-h", "--help"):
            args["help"] = True
        elif a == "--self-test":
            args["self_test"] = True
        elif a == "--json":
            args["json"] = True
        elif a == "--output" and i + 1 < len(argv):
            args["output"] = argv[i + 1]; i += 1
        elif a.startswith("--output="):
            args["output"] = a[len("--output="):]
        elif a == "--outdir" and i + 1 < len(argv):
            args["outdir"] = argv[i + 1]; i += 1
        elif a.startswith("--outdir="):
            args["outdir"] = a[len("--outdir="):]
        elif a == "--engine" and i + 1 < len(argv):
            args["engine"] = argv[i + 1]; i += 1
        elif a.startswith("--engine="):
            args["engine"] = a[len("--engine="):]
        elif not a.startswith("-"):
            args["inputs"].append(a)
        i += 1
    return args


def output_path_for(in_abs, args):
    if args["output"]:
        return to_windows_abs(args["output"])
    base = Path(in_abs).stem + ".pdf"
    if args["outdir"]:
        return os.path.join(to_windows_abs(args["outdir"]), base)
    return os.path.join(os.path.dirname(in_abs), base)


def self_test():
    try:
        from pptx import Presentation
    except ImportError:
        sys.stderr.write(
            "missing dep: python-pptx. Run via `uv run convert-pptx-to-pdf.py --self-test`.\n"
        )
        sys.exit(1)

    tmpdir = tempfile.mkdtemp(prefix="pptx2pdf_selftest_")
    try:
        in_path = os.path.join(tmpdir, "selftest.pptx")
        out_path = os.path.join(tmpdir, "selftest.pdf")

        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[0])  # title slide
        s1.shapes.title.text = "Self Test Slide 1"
        try:
            s1.placeholders[1].text = "pptx2pdf selftest marker"
        except (KeyError, IndexError):
            pass
        s2 = prs.slides.add_slide(prs.slide_layouts[5])   # title only
        s2.shapes.title.text = "Self Test Slide 2"
        prs.save(in_path)

        try:
            engine_name, engine_path = resolve_engine("auto")
        except RuntimeError as e:
            sys.stderr.write(f"selftest SKIP/FAIL: {e}\n")
            sys.exit(1)

        convert(engine_name, engine_path, to_windows_abs(in_path), to_windows_abs(out_path))

        if not os.path.isfile(out_path):
            sys.stderr.write("selftest FAIL: no PDF produced\n")
            sys.exit(1)
        size = os.path.getsize(out_path)
        if size < 1000:
            sys.stderr.write(f"selftest FAIL: PDF suspiciously small ({size} bytes)\n")
            sys.exit(1)
        pages = pdf_page_count(out_path)
        if pages is not None and pages != 2:
            sys.stderr.write(f"selftest FAIL: expected 2 pages, got {pages}\n")
            sys.exit(1)
        sys.stdout.write(f"selftest pass (engine={engine_name}, {size} bytes, pages={pages})\n")
        sys.exit(0)
    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)


def main():
    args = parse_args(sys.argv[1:])
    if args["help"]:
        sys.stdout.write(__doc__)
        return
    if args["self_test"]:
        return self_test()
    if not args["inputs"]:
        sys.stderr.write(USAGE)
        sys.exit(2)
    if args["output"] and len(args["inputs"]) > 1:
        sys.stderr.write("--output takes a single input; use --outdir for multiple files.\n")
        sys.exit(2)

    norm_inputs = []
    for p in args["inputs"]:
        ab = to_windows_abs(p)
        if not os.path.isfile(ab):
            sys.stderr.write(f"file not found: {p}\n")
            sys.exit(1)
        if Path(ab).suffix.lower() not in PPT_EXTS:
            sys.stderr.write(f"not a PowerPoint file ({', '.join(PPT_EXTS)}): {p}\n")
            sys.exit(1)
        norm_inputs.append(ab)

    try:
        engine_name, engine_path = resolve_engine(args["engine"])
    except RuntimeError as e:
        sys.stderr.write(str(e) + "\n")
        sys.exit(1)

    results = []
    had_error = False
    for in_abs in norm_inputs:
        out_abs = output_path_for(in_abs, args)
        os.makedirs(os.path.dirname(out_abs) or ".", exist_ok=True)
        try:
            convert(engine_name, engine_path, in_abs, out_abs)
            results.append({
                "input": in_abs, "output": out_abs, "engine": engine_name,
                "pages": pdf_page_count(out_abs), "bytes": os.path.getsize(out_abs),
                "status": "success",
            })
        except Exception as e:
            had_error = True
            results.append({
                "input": in_abs, "output": out_abs, "engine": engine_name,
                "pages": None, "bytes": None, "status": "error",
                "error": f"{type(e).__name__}: {e}",
            })

    if args["json"]:
        sys.stdout.write(json.dumps({"results": results}, indent=2) + "\n")
    else:
        for r in results:
            if r["status"] == "success":
                sys.stdout.write(r["output"] + "\n")
            else:
                sys.stderr.write(f"FAILED {r['input']}: {r['error']}\n")

    if had_error:
        sys.exit(1)


if __name__ == "__main__":
    main()
