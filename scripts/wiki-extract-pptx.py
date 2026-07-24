#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0"]
# ///
"""
wiki-extract-pptx.py — extract markdown from a PowerPoint deck via python-pptx.

Stage 12 of the wiki upgrade arc — see [[Eric Ma — Mastering PKM with Obsidian
and AI]] for the canonical recipe.

Usage:
  uv run wiki-extract-pptx.py <path-to-pptx>      # markdown to stdout
  uv run wiki-extract-pptx.py --self-test          # synthesize tmp PPTX, parse, exit 0/1

LIMITATIONS:
- Charts, embedded images, SmartArt, animations NOT extracted (text + tables only).
  Eric Ma's recipe pairs python-pptx with LibreOffice + PIL + VLM captioning
  for ~90-95% accuracy. That dual path is out of scope here.
- Tier classification can mark graphics-heavy decks as `secondary` or
  `reference` if missing content is load-bearing.

First run cost: ~30-60 sec (uv resolves python-pptx + lxml wheels). Cached after.
"""
import sys
import tempfile
from pathlib import Path

# Force UTF-8 on stdout/stderr so non-ASCII content (accented chars, em-dashes,
# math symbols) doesn't crash on Windows where the default codec is cp1252.
# Python 3.7+ feature; harmless on Linux/macOS.
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except (AttributeError, ValueError):
    pass

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write(
        "missing dep: python-pptx. Run via `uv run wiki-extract-pptx.py ...` "
        "(PEP-723 inline deps), not bare python.\n"
    )
    sys.exit(1)


def _shape_text(shape):
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _table_md(shape):
    if not shape.has_table:
        return ""
    tbl = shape.table
    rows = []
    for row in tbl.rows:
        cells = [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    cols = max(len(r) for r in rows)
    rows = [r + [""] * (cols - len(r)) for r in rows]
    head = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * cols) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join([head, sep] + body)


def slide_to_markdown(slide, idx):
    out = [f"## Slide {idx}"]
    title = ""
    body_lines = []
    tables = []

    for shape in slide.shapes:
        # Title placeholder is idx 0 by convention
        if shape.has_text_frame and getattr(shape, "is_placeholder", False):
            try:
                if shape.placeholder_format.idx == 0 and not title:
                    title = _shape_text(shape).split("\n", 1)[0].strip()
                    continue
            except (AttributeError, ValueError):
                pass
        if shape.has_text_frame:
            for line in _shape_text(shape).split("\n"):
                if line.strip():
                    body_lines.append(line.strip())
        elif getattr(shape, "has_table", False):
            md = _table_md(shape)
            if md:
                tables.append(md)

    if title:
        out.append(f"### {title}")
    for line in body_lines:
        out.append(f"- {line}")
    for tbl in tables:
        out.append("")
        out.append(tbl)

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            out.append("")
            out.append(f"> **Notes:** {notes}")

    return "\n".join(out)


def self_test():
    KNOWN = "Stage 12 pptx selftest marker"
    KNOWN_TITLE = "Self Test Title"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp_path = Path(f.name)
    try:
        prs = Presentation()
        layout = prs.slide_layouts[1]  # title + content
        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = KNOWN_TITLE
        body_ph = slide1.placeholders[1].text_frame
        body_ph.text = KNOWN
        body_ph.add_paragraph().text = "second bullet"

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        slide2.shapes.title.text = "Slide Two"

        prs.save(str(tmp_path))

        rt = Presentation(str(tmp_path))
        out_md = []
        for i, slide in enumerate(rt.slides, 1):
            out_md.append(slide_to_markdown(slide, i))
        md = "\n".join(out_md)

        if KNOWN not in md:
            sys.stderr.write(f"selftest FAIL: marker not in extracted output\n  got:\n{md}\n")
            sys.exit(1)
        if f"### {KNOWN_TITLE}" not in md:
            sys.stderr.write(f"selftest FAIL: title not detected\n  got:\n{md}\n")
            sys.exit(1)
        if "## Slide 1" not in md or "## Slide 2" not in md:
            sys.stderr.write(f"selftest FAIL: both slides expected\n  got:\n{md}\n")
            sys.exit(1)
        sys.stdout.write("selftest pass\n")
        sys.exit(0)
    finally:
        try:
            tmp_path.unlink()
        except OSError:
            pass


def main():
    if len(sys.argv) >= 2 and sys.argv[1] == "--self-test":
        return self_test()
    if len(sys.argv) < 2:
        sys.stderr.write("usage: uv run wiki-extract-pptx.py <path-to-pptx> | --self-test\n")
        sys.exit(2)

    path = Path(sys.argv[1])
    if not path.exists():
        sys.stderr.write(f"file not found: {path}\n")
        sys.exit(1)

    try:
        prs = Presentation(str(path))
    except Exception as e:
        sys.stderr.write(f"python-pptx failed: {type(e).__name__}: {e}\n")
        sys.exit(1)

    out = [
        f"# {path.stem}",
        "",
        f"_Source: `{path.name}` · {len(prs.slides)} slides_",
        "",
    ]
    for i, slide in enumerate(prs.slides, 1):
        out.append(slide_to_markdown(slide, i))
        out.append("")
    sys.stdout.write("\n".join(out))


if __name__ == "__main__":
    main()
